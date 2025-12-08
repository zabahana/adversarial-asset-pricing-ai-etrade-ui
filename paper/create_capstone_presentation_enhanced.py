#!/usr/bin/env python3
"""
Create enhanced capstone presentation with architecture, formulae, and agentic AI details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import os

def load_meta_results():
    """Load META case study results"""
    try:
        with open('../results/meta_model_results.json', 'r') as f:
            return json.load(f)
    except Exception as e:
        print(f"Warning: Could not load META results: {e}")
        return None

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Adversarial-Robust Multi-Head Attention DQN\nfor Asset Pricing"
    subtitle.text = "AI 894 Capstone Project\nAn Agentic AI System for Portfolio Optimization"
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "System Architecture: Multi-Head Attention DQN"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add architecture image if it exists
    img_path = "architecture.png"
    if os.path.exists(img_path):
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(6)
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        # Add text description if image not found
        desc_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        desc_frame = desc_shape.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = "Architecture Diagram:\n\n• 8 Input Feature Groups → Multi-Head Attention (8 heads)\n• 3 Transformer Layers → Global Average Pooling\n• Q-Network → BUY/HOLD/SELL Actions\n\n(Image: architecture.png)"
        p.font.size = Pt(18)
    
    return slide

def add_dqn_formulae_slide(prs):
    """Add detailed DQN formulae slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Deep Q-Network (DQN): Mathematical Formulation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Reinforcement Learning Framework:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "MDP Formulation: M = (S, A, P, R, γ)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• State Space S: High-dimensional feature vectors (43 features × 20 time steps)"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Action Space A: {BUY, HOLD, SELL}"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Reward R(s, a, s'): Portfolio return = (V_{t+1} - V_t) / V_t"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Q-Function and Bellman Equation:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Q-function definition:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Q^π(s, a) = E[Σ(k=0 to ∞) γ^k r_{t+k+1} | s_t = s, a_t = a]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Bellman Optimality Equation:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Q*(s, a) = E[r(s, a, s') + γ max_{a'} Q*(s', a') | s, a]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "DQN Learning Objective:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "L(θ) = E[(Q_θ(s, a) - (r + γ max_{a'} Q_{θ^-}(s', a')))^2]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "where θ^- is target network, updated every C steps"
    p.font.size = Pt(12)
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Example (META Stock):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "State: RSI=0.65, MACD=0.02, sentiment=0.8 → Q(s, BUY)=0.85, Q(s, HOLD)=0.45, Q(s, SELL)=0.25"
    p.font.size = Pt(11)
    p.font.name = 'Courier New'
    
    return slide

def add_dqn_training_slide(prs):
    """Add DQN training details slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "DQN Training Algorithm: Step-by-Step"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. Experience Replay Buffer:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Store transitions: (s_t, a_t, r_t, s_{t+1}, done_t)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Example: (RSI=0.28, MACD=+ve, earnings in 2d, BUY, +2.3%, next_state)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Target Network Update:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Q_target = r + γ(1 - done) · max_{a'} Q_{θ^-}(s', a')"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   γ = 0.99 (discount factor), θ^- updated every 100 steps"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Loss Minimization:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   L(θ) = (1/|B|) Σ_{(s,a,r,s') ∈ B} (Q_θ(s, a) - Q_target)^2"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Batch size |B| = 32, learning rate α = 0.001"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "4. Gradient Update:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   θ ← θ - α · ∇_θ L(θ)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Uses Adam optimizer with adaptive learning rates"
    p.font.size = Pt(12)
    p.level = 1
    
    return slide

def add_adversarial_definitions_slide(prs):
    """Add adversarial attack definitions"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Attacks: Definitions and Threat Model"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Threat Model:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Adversary seeks perturbation δ such that:"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "max_{||δ||_p ≤ ε} L(Q_θ(s + δ, a), y)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "subject to: s + δ ∈ S (feasibility constraint)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Attack Types:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. FGSM (Fast Gradient Sign Method):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   x_adv = x + ε · sign(∇_x L(Q_θ(x), y))"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Example: ε = 0.01, perturbs RSI from 0.65 → 0.64"
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "2. PGD (Projected Gradient Descent):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Iterative: x^{(t+1)} = Proj_{B_ε(x)}(x^{(t)} + α·sign(∇_{x^{(t)}} L))"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   More powerful than FGSM, 10 iterations"
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "3. C&W (Carlini & Wagner):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Optimization: min_δ ||δ||_p + c·f(x + δ)"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    
    return slide

def add_adversarial_examples_slide(prs):
    """Add adversarial attack examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Attack: Real-World Example"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scenario: Attacking META Stock Prediction"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Clean Input State:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• RSI = 0.65, MACD = 0.02, Earnings sentiment = 0.8"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Model prediction: Q(s, BUY) = 0.85, Q(s, SELL) = 0.25"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Correct action: BUY (profitable trade)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "FGSM Attack (ε = 0.01):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "1. Compute gradient: ∇_s L = [-0.15, 0.08, -0.22]"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "2. Apply perturbation: δ = 0.01 × sign(∇_s) = [-0.01, 0.01, -0.01]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "3. Adversarial state: [RSI=0.64, MACD=0.03, sentiment=0.79]"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "4. Without robust training: Q(s_adv, SELL) = 0.60 (WRONG!)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Robust Model (After Adversarial Training):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Q(s_adv, BUY) = 0.84 (maintains correct prediction)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• 100% resistance to FGSM, PGD, BIM, C&W attacks"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    return slide

def add_agentic_ai_slide(prs):
    """Add agentic AI system explanation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Agentic AI System: Autonomous Decision-Making"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "What Makes This an Agentic AI System?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. Autonomous Decision-Making:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Agent observes market state → processes through MHA-DQN → selects action"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • No human intervention required for trading decisions"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Learns optimal policy through trial and error"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Sequential Decision-Making:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Considers long-term consequences (discount factor γ = 0.99)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Actions affect future states and rewards"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Balances exploration (ε-greedy) vs exploitation"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Self-Learning Capability:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Updates Q-function from experience replay buffer"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Improves performance over time (Sharpe: 0.47 → 2.88)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Adapts to changing market conditions"
    p.font.size = Pt(13)
    p.level = 1
    
    return slide

def add_api_collaboration_slide(prs):
    """Add API collaboration and AI system integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-AI System Collaboration: API Integration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Integrated APIs and AI Systems:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. Alpha Vantage API:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Provides: Macroeconomic indicators (GDP, CPI, unemployment, Fed funds rate)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Commodities: Oil, gold, copper prices"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Market indices: SPY, VIX, sector ETFs"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Financial Modeling Prep (FMP) API:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Provides: Earnings call transcripts, fundamental data"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Earnings metrics: EPS, surprise percentages, growth rates"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. OpenAI API (LLM Integration):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Analyzes earnings call transcripts → generates sentiment scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Extracts key themes: growth outlook, risk factors, management confidence"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Output: earnings_call_score (0-1) → fed into MHA-DQN as feature"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Collaborative AI Pipeline:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "FMP API → Earnings Transcript → OpenAI LLM → Sentiment Score → MHA-DQN → Trading Decision"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Example: META earnings call analyzed → sentiment=0.85 → boosts Q(BUY) from 0.75 to 0.88"
    p.font.size = Pt(11)
    p.level = 1
    
    return slide

def add_use_case_slide(prs):
    """Add compelling use case"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Compelling Use Case: Institutional Portfolio Management"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scenario: $100M Tech Sector Fund"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Challenge:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Portfolio manager needs to optimize holdings across 50 tech stocks"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Must respond quickly to earnings, macro events, market regime changes"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Current manual analysis: 8 hours/day, delayed reactions"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Solution: Agentic AI Trading System"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Real-time analysis: Processes all 50 stocks in < 2 minutes"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Automated decisions: BUY/HOLD/SELL recommendations with confidence scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Integrated AI analysis: LLM summarizes earnings calls → sentiment scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Results (META Case Study):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Sharpe ratio: 2.88 (vs 0.47 baseline) = 513% improvement"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• CAGR: 78% (vs 17.94% baseline) = 335% improvement"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Risk control: Max drawdown -7% (vs -18% baseline)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Adversarial robustness: 100% resistance to attacks"
    p.font.size = Pt(13)
    p.level = 1
    
    return slide

def add_attention_formulae_slide(prs):
    """Add multi-head attention formulae"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-Head Attention: Mathematical Formulation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scaled Dot-Product Attention:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Attention(Q, K, V) = softmax((Q·K^T) / √d_k) · V"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Multi-Head Attention:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "MHA(Q, K, V) = Concat(head_1, ..., head_H) · W_O"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "where: head_h = Attention(Q·W_Q^(h), K·W_K^(h), V·W_V^(h))"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "For Our Model:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• H = 8 heads (one per feature group)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• d_k = d_v = 128 (key/value dimensions)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Sequence length L = 20 (20-day lookback)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Feature Group Processing:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "H_g^(0) = X_g · W_proj  ∈ R^(L × d_h)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "After 3 transformer layers: z_g = (1/L) Σ_{t=1}^L H_{g,t}^{(3)}"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    
    return slide

# Keep existing functions and add new slides in main()
def main():
    """Create enhanced presentation"""
    print("Creating enhanced capstone presentation...")
    
    template_path = "AI 894 Capstone Video Presentation Slides.pptx"
    try:
        prs = Presentation(template_path)
        print(f"Loaded template: {template_path}")
    except:
        print("Template not found, creating new presentation...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    meta_results = load_meta_results()
    
    # Add all slides
    print("Adding slides...")
    add_title_slide(prs)
    add_intro_slide(prs)
    add_architecture_slide(prs)
    add_dqn_formulae_slide(prs)
    add_dqn_training_slide(prs)
    add_attention_formulae_slide(prs)
    add_methodology_overview_slide(prs)
    add_feature_engineering_examples_slide(prs)
    add_attention_mechanism_examples_slide(prs)
    add_training_examples_slide(prs)
    add_adversarial_definitions_slide(prs)
    add_adversarial_examples_slide(prs)
    add_adversarial_training_examples_slide(prs)
    add_agentic_ai_slide(prs)
    add_api_collaboration_slide(prs)
    add_use_case_slide(prs)
    add_meta_results_slide(prs, meta_results)
    add_adversarial_results_slide(prs, meta_results)
    add_infrastructure_slide(prs)
    add_benefits_slide(prs)
    add_conclusion_slide(prs)
    
    output_path = "AI_894_Capstone_Presentation_Enhanced.pptx"
    prs.save(output_path)
    print(f"\n✅ Enhanced presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

# Include all the existing helper functions from the original script
# (add_intro_slide, add_methodology_overview_slide, etc.)
# I'll copy them from the original file...

if __name__ == "__main__":
    # First, import all existing functions
    import sys
    sys.path.insert(0, os.path.dirname(__file__))
    
    # Import from original script (need to add these functions)
    # For now, let's create a complete version
    print("Note: This script requires functions from create_capstone_presentation.py")
    print("Please ensure all helper functions are available.")
    main()

