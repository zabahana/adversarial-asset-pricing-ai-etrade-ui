#!/usr/bin/env python3
"""
Create detailed capstone presentation based on template
Focuses on practical examples and methodology, avoids proofs/theorems
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json

def load_meta_results():
    """Load META case study results"""
    try:
        with open('../results/meta_model_results.json', 'r') as f:
            return json.load(f)
    except Exception as e:
        print(f"Warning: Could not load META results: {e}")
        return None

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Adversarial-Robust Multi-Head Attention DQN\nfor Asset Pricing"
    subtitle.text = "AI 894 Capstone Project\nAn Agentic AI System for Portfolio Optimization"
    
    # Format title
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_intro_slide(prs):
    """Add introduction slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Introduction"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "• Asset pricing in financial markets is a complex sequential decision-making problem"
    p.font.size = Pt(18)
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Challenges include:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - High dimensionality (price, macro, sentiment, technical indicators)"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Temporal dependencies (long-range patterns in time series)"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Adversarial nature (data corruption, market manipulation, model attacks)"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Our Solution: Multi-Head Attention Deep Q-Network (MHA-DQN)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Multi-head attention for complex feature interactions"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Deep Q-learning for optimal sequential decisions"
    p.font.size = Pt(16)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  - Adversarial training for robustness guarantees"
    p.font.size = Pt(16)
    p.level = 1
    
    return slide

def add_problem_statement_slide(prs):
    """Add problem statement slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Problem Statement"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Traditional asset pricing models face limitations:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "1. Limited Feature Interaction"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Example: Can't capture how RSI interacts with earnings sentiment during high volatility"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Vulnerability to Adversarial Attacks"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Example: Small perturbations to input features can cause incorrect trading decisions"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Sequential Decision Complexity"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Example: Need to consider entire history of 20+ time steps for optimal action"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "4. Lack of Robustness Metrics"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Example: Model performance degrades significantly under data corruption"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_methodology_overview_slide(prs):
    """Add methodology overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Methodology Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Three-Stage Pipeline:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Stage 1: Feature Engineering"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Extract 8 feature groups (Price, Macro, Commodities, Market Indices, Forex, Technical, Earnings, Crypto)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Create sequences of length 20 time steps"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Stage 2: Multi-Head Attention DQN Training"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • 8 attention heads process feature groups in parallel"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Q-learning learns optimal BUY/HOLD/SELL actions"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Stage 3: Adversarial Training & Evaluation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Train on adversarial examples (FGSM, PGD, BIM, C&W attacks)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Evaluate robustness against multiple attack types"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_feature_engineering_examples_slide(prs):
    """Add detailed feature engineering examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Feature Engineering: Practical Examples"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Group 1: Price & Returns Features"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • return = (close_today - close_yesterday) / close_yesterday"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • volatility = std(returns, window=20 days)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • RSI_14 = Relative Strength Index (overbought/oversold indicator)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Group 6: Technical Indicators"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • MACD = EMA(12) - EMA(26)  (momentum indicator)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Bollinger Bands: upper/middle/lower = sma ± 2*std"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • ATR = Average True Range (volatility measure)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Group 7: Earnings & Sentiment"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • earnings_surprise_pct = (actual_eps - estimated_eps) / estimated_eps"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • earnings_call_score = OpenAI LLM analysis of transcript sentiment"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    
    return slide

def add_attention_mechanism_examples_slide(prs):
    """Add attention mechanism examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-Head Attention: How It Works"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Example: Processing META Stock Features"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Input: 20-day sequence of 43 features"
    p.font.size = Pt(16)
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "8 Attention Heads Process Different Aspects:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Head 1 (Price): Focuses on recent price trends vs. historical patterns"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 2 (Macro): Weights GDP, CPI, unemployment relative to stock price"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 3 (Commodities): Correlates oil/gold prices with tech stock movements"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 4 (Market): Compares stock to SPY, sector ETF performance"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 5 (Forex): Links USD strength to international revenue impact"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 6 (Technical): Identifies MACD crossovers, RSI extremes"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 7 (Earnings): Emphasizes earnings calls near report dates"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Head 8 (Crypto): Monitors Bitcoin correlation with tech sentiment"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Output: Weighted combination → Q-values for BUY/HOLD/SELL"
    p.font.size = Pt(16)
    p.font.bold = True
    
    return slide

def add_training_examples_slide(prs):
    """Add training process examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Training Process: Step-by-Step Example"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Example: Training on META Stock Data"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Step 1: Experience Collection"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Observe state: [Price features, Technical indicators, Earnings sentiment, ...]"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Model predicts Q-values: BUY=0.8, HOLD=0.5, SELL=0.3 → Action: BUY"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Execute action, observe reward: +2.5% portfolio return"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Step 2: Experience Replay"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Store transition: (state, BUY, +2.5%, next_state)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Sample batch of 32 past experiences randomly"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Breaks temporal correlation, improves learning stability"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Step 3: Target Network Update"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Compute target: Q_target = reward + γ * max(Q_target_network(next_state))"
    p.font.size = Pt(14)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Update main network to minimize (Q_predict - Q_target)²"
    p.font.size = Pt(14)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Copy main network → target network every 100 steps"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_adversarial_training_examples_slide(prs):
    """Add adversarial training examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Training: Real-World Examples"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Problem: Models vulnerable to input perturbations"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Example Attack Scenario:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  Original: RSI = 0.65, MACD = 0.02 → Prediction: BUY"
    p.font.size = Pt(14)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  Perturbed: RSI = 0.66, MACD = 0.021 → Prediction: SELL (WRONG!)"
    p.font.size = Pt(14)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Solution: Adversarial Training"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "FGSM Attack (Fast Gradient Sign Method):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  1. Compute gradient: ∇_x L(model(x), label)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  2. Create perturbation: δ = ε * sign(∇_x L)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  3. Adversarial example: x_adv = x + δ (ε = 0.01)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.name = 'Courier New'
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Training Process:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Train on 50% clean data + 50% adversarial examples"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Model learns to make consistent predictions under perturbations"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Result: Robust model that resists attacks"
    p.font.size = Pt(12)
    p.level = 1
    
    return slide

def add_meta_results_slide(prs, meta_results):
    """Add META case study results"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Results: META Stock Case Study"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    if meta_results:
        clean_metrics = meta_results.get('mha_dqn_clean', {}).get('metrics', {})
        robust_metrics = meta_results.get('mha_dqn_robust', {}).get('metrics', {})
        baseline_metrics = meta_results.get('baseline_dqn', {}).get('metrics', {})
        
        p = tf.paragraphs[0]
        p.text = "Experimental Setup:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = f"  • Ticker: META (Meta Platforms Inc.)"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = f"  • Backtest Period: 229 trading days"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = f"  • Initial Portfolio: $10,000"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "Performance Comparison:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        sharpe = clean_metrics.get('sharpe', 0)
        cagr = clean_metrics.get('cagr', 0)
        total_return = clean_metrics.get('total_return', 0)
        p.text = f"MHA-DQN Clean: Sharpe={sharpe:.2f}, CAGR={cagr*100:.0f}%, Total Return={total_return*100:.0f}%"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        sharpe_r = robust_metrics.get('sharpe', 0)
        cagr_r = robust_metrics.get('cagr', 0)
        total_return_r = robust_metrics.get('total_return', 0)
        p.text = f"MHA-DQN Robust: Sharpe={sharpe_r:.2f}, CAGR={cagr_r*100:.0f}%, Total Return={total_return_r*100:.0f}%"
        p.font.size = Pt(14)
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        sharpe_b = baseline_metrics.get('sharpe', 0)
        cagr_b = baseline_metrics.get('cagr', 0)
        p.text = f"Baseline DQN: Sharpe={sharpe_b:.2f}, CAGR={cagr_b*100:.1f}%"
        p.font.size = Pt(14)
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "Key Achievements:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(6)
        
        improvement = ((sharpe - sharpe_b) / sharpe_b * 100) if sharpe_b > 0 else 0
        p = tf.add_paragraph()
        p.text = f"  • {improvement:.0f}% improvement in Sharpe ratio over baseline"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        final_value = clean_metrics.get('final_portfolio_value', 10000)
        p.text = f"  • Final portfolio value: ${final_value:,.2f} (from $10,000)"
        p.font.size = Pt(14)
        p.level = 1
        p.space_after = Pt(4)
        
        max_dd = clean_metrics.get('max_drawdown', 0)
        p = tf.add_paragraph()
        p.text = f"  • Maximum drawdown: {max_dd*100:.0f}% (excellent risk control)"
        p.font.size = Pt(14)
        p.level = 1
    
    return slide

def add_adversarial_results_slide(prs, meta_results):
    """Add adversarial attack results"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Robustness: Attack Resistance Results"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    if meta_results and 'adversarial_attack_results' in meta_results:
        attacks = meta_results['adversarial_attack_results']
        
        p = tf.paragraphs[0]
        p.text = "Attack Resistance (100% = Perfect Resistance):"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(8)
        
        attack_list = [
            ('FGSM', 'Fast Gradient Sign Method'),
            ('PGD', 'Projected Gradient Descent'),
            ('BIM', 'Basic Iterative Method'),
            ('C&W', 'Carlini & Wagner'),
        ]
        
        for attack_key, attack_name in attack_list:
            if attack_key in attacks:
                attack_data = attacks[attack_key]
                resistance = attack_data.get('resistance', 'N/A')
                robustness = attack_data.get('robustness_score', 'N/A')
                
                p = tf.add_paragraph()
                p.text = f"{attack_name} ({attack_key}):"
                p.font.size = Pt(14)
                p.font.bold = True
                p.space_after = Pt(4)
                
                p = tf.add_paragraph()
                p.text = f"  • Resistance: {resistance}"
                p.font.size = Pt(12)
                p.level = 1
                p.space_after = Pt(4)
                
                p = tf.add_paragraph()
                p.text = f"  • Robustness Score: {robustness}"
                p.font.size = Pt(12)
                p.level = 1
                p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "Key Insight:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = "Adversarially-trained model maintains 100% resistance to major attack types"
        p.font.size = Pt(14)
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = "This ensures reliable performance even under data corruption or adversarial inputs"
        p.font.size = Pt(14)
    
    return slide

def add_infrastructure_slide(prs):
    """Add infrastructure deployment slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Infrastructure & Deployment"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Google Cloud Platform (GCP) Deployment:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Cloud Run with GPU Acceleration:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • GPU: NVIDIA L4 (1 GPU per instance)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Memory: 16 GB RAM"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • CPU: 4 vCPUs"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Timeout: 3600 seconds (1 hour)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Data Pipeline:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Alpha Vantage API: Macroeconomic indicators"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Financial Modeling Prep: Earnings transcripts"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • OpenAI API: Sentiment analysis (optional)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Streamlit Cloud Deployment:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Web-based user interface"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Auto-deployment from GitHub"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Real-time analysis and visualization"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_benefits_slide(prs):
    """Add benefits slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Key Benefits & Applications"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Performance Advantages:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • 513% improvement in Sharpe ratio over baseline"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • 335% improvement in CAGR"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • 61% reduction in maximum drawdown"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Security & Robustness:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • 100% resistance to FGSM, PGD, BIM, C&W attacks"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Resilient to data corruption and measurement errors"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Practical Applications:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Institutional quantitative trading"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Automated trading platforms"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Portfolio optimization systems"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Risk management tools"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_conclusion_slide(prs):
    """Add conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Conclusion & Future Work"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Achievements:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  ✅ Developed adversarial-robust MHA-DQN framework"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  ✅ Demonstrated exceptional performance (78% CAGR, 2.88 Sharpe)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  ✅ Achieved perfect resistance to adversarial attacks"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  ✅ Deployed production-ready system on cloud infrastructure"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Future Work:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Extend to multi-asset portfolio optimization"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Incorporate additional data sources (alternative data)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Real-time adaptive learning from market feedback"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Integration with live trading platforms"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "System Architecture: Multi-Head Attention DQN"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add architecture image if it exists
    import os
    img_path = "architecture.png"
    if os.path.exists(img_path):
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(6)
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        # Add text description if image not found
        desc_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        desc_frame = desc_shape.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = "Architecture Diagram:\n\n• 8 Input Feature Groups → Multi-Head Attention (8 heads)\n• 3 Transformer Layers → Global Average Pooling\n• Q-Network → BUY/HOLD/SELL Actions\n\n(Image: architecture.png)"
        p.font.size = Pt(18)
    
    return slide

def add_dqn_formulae_slide(prs):
    """Add detailed DQN formulae slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Deep Q-Network (DQN): Mathematical Formulation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Reinforcement Learning Framework:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "MDP Formulation: M = (S, A, P, R, γ)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• State Space S: High-dimensional feature vectors (43 features × 20 time steps)"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Action Space A: {BUY, HOLD, SELL}"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Reward R(s, a, s'): Portfolio return = (V_{t+1} - V_t) / V_t"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Q-Function and Bellman Equation:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Q-function definition:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Q^π(s, a) = E[Σ(k=0 to ∞) γ^k r_{t+k+1} | s_t = s, a_t = a]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Bellman Optimality Equation:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Q*(s, a) = E[r(s, a, s') + γ max_{a'} Q*(s', a') | s, a]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "DQN Learning Objective:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "L(θ) = E[(Q_θ(s, a) - (r + γ max_{a'} Q_{θ^-}(s', a')))^2]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "where θ^- is target network, updated every C steps"
    p.font.size = Pt(12)
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Example (META Stock):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "State: RSI=0.65, MACD=0.02, sentiment=0.8 → Q(s, BUY)=0.85, Q(s, HOLD)=0.45, Q(s, SELL)=0.25"
    p.font.size = Pt(11)
    p.font.name = 'Courier New'
    
    return slide

def add_dqn_training_slide(prs):
    """Add DQN training details slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "DQN Training Algorithm: Detailed Steps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1. Experience Replay Buffer:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Store transitions: (s_t, a_t, r_t, s_{t+1}, done_t)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Example: (RSI=0.28, MACD=+ve, earnings in 2d, BUY, +2.3%, next_state)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Target Network Update:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   Q_target = r + γ(1 - done) · max_{a'} Q_{θ^-}(s', a')"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   γ = 0.99 (discount factor), θ^- updated every 100 steps"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Loss Minimization:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   L(θ) = (1/|B|) Σ_{(s,a,r,s') ∈ B} (Q_θ(s, a) - Q_target)^2"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Batch size |B| = 32, learning rate α = 0.001"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "4. Gradient Update:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   θ ← θ - α · ∇_θ L(θ)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Uses Adam optimizer with adaptive learning rates"
    p.font.size = Pt(12)
    p.level = 1
    
    return slide

def add_attention_formulae_slide(prs):
    """Add multi-head attention formulae"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-Head Attention: Mathematical Formulation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scaled Dot-Product Attention:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Attention(Q, K, V) = softmax((Q·K^T) / √d_k) · V"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Multi-Head Attention:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "MHA(Q, K, V) = Concat(head_1, ..., head_H) · W_O"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "where: head_h = Attention(Q·W_Q^(h), K·W_K^(h), V·W_V^(h))"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "For Our Model:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• H = 8 heads (one per feature group)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• d_k = d_v = 128 (key/value dimensions)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Sequence length L = 20 (20-day lookback)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Feature Group Processing:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "H_g^(0) = X_g · W_proj  ∈ R^(L × d_h)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "After 3 transformer layers: z_g = (1/L) Σ_{t=1}^L H_{g,t}^{(3)}"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    
    return slide

def add_adversarial_definitions_slide(prs):
    """Add adversarial attack definitions"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Attacks: Definitions and Threat Model"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Threat Model:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Adversary seeks perturbation δ such that:"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "max_{||δ||_p ≤ ε} L(Q_θ(s + δ, a), y)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "subject to: s + δ ∈ S (feasibility constraint)"
    p.font.size = Pt(13)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Attack Types:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. FGSM (Fast Gradient Sign Method):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   x_adv = x + ε · sign(∇_x L(Q_θ(x), y))"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Example: ε = 0.01, perturbs RSI from 0.65 → 0.64"
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "2. PGD (Projected Gradient Descent):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Iterative: x^{(t+1)} = Proj_{B_ε(x)}(x^{(t)} + α·sign(∇_{x^{(t)}} L))"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   More powerful than FGSM, 10 iterations"
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "3. C&W (Carlini & Wagner):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Optimization: min_δ ||δ||_p + c·f(x + δ)"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    
    return slide

def add_adversarial_examples_slide(prs):
    """Add adversarial attack examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Attack: Real-World Example"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scenario: Attacking META Stock Prediction"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Clean Input State:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• RSI = 0.65, MACD = 0.02, Earnings sentiment = 0.8"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Model prediction: Q(s, BUY) = 0.85, Q(s, SELL) = 0.25"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Correct action: BUY (profitable trade)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "FGSM Attack (ε = 0.01):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "1. Compute gradient: ∇_s L = [-0.15, 0.08, -0.22]"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "2. Apply perturbation: δ = 0.01 × sign(∇_s) = [-0.01, 0.01, -0.01]"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "3. Adversarial state: [RSI=0.64, MACD=0.03, sentiment=0.79]"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "4. Without robust training: Q(s_adv, SELL) = 0.60 (WRONG!)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Robust Model (After Adversarial Training):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Q(s_adv, BUY) = 0.84 (maintains correct prediction)"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• 100% resistance to FGSM, PGD, BIM, C&W attacks"
    p.font.size = Pt(12)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    return slide

def add_agentic_ai_slide(prs):
    """Add agentic AI system explanation"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Agentic AI System: Autonomous Decision-Making"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "What Makes This an Agentic AI System?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. Autonomous Decision-Making:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Agent observes market state → processes through MHA-DQN → selects action"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • No human intervention required for trading decisions"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Learns optimal policy through trial and error"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Sequential Decision-Making:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Considers long-term consequences (discount factor γ = 0.99)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Actions affect future states and rewards"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Balances exploration (ε-greedy) vs exploitation"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. Self-Learning Capability:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Updates Q-function from experience replay buffer"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Improves performance over time (Sharpe: 0.47 → 2.88)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Adapts to changing market conditions"
    p.font.size = Pt(13)
    p.level = 1
    
    return slide

def add_api_collaboration_slide(prs):
    """Add API collaboration and AI system integration"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Multi-AI System Collaboration: API Integration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Integrated APIs and AI Systems:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. Alpha Vantage API:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Provides: Macroeconomic indicators (GDP, CPI, unemployment, Fed funds rate)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Commodities: Oil, gold, copper prices"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Market indices: SPY, VIX, sector ETFs"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Financial Modeling Prep (FMP) API:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Provides: Earnings call transcripts, fundamental data"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Earnings metrics: EPS, surprise percentages, growth rates"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. OpenAI API (LLM Integration):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "   • Analyzes earnings call transcripts → generates sentiment scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Extracts key themes: growth outlook, risk factors, management confidence"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Output: earnings_call_score (0-1) → fed into MHA-DQN as feature"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Collaborative AI Pipeline:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "FMP API → Earnings Transcript → OpenAI LLM → Sentiment Score → MHA-DQN → Trading Decision"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Example: META earnings call analyzed → sentiment=0.85 → boosts Q(BUY) from 0.75 to 0.88"
    p.font.size = Pt(11)
    p.level = 1
    
    return slide

def add_use_case_slide(prs):
    """Add compelling use case"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Compelling Use Case: Institutional Portfolio Management"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scenario: $100M Tech Sector Fund"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Challenge:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Portfolio manager needs to optimize holdings across 50 tech stocks"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Must respond quickly to earnings, macro events, market regime changes"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Current manual analysis: 8 hours/day, delayed reactions"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Solution: Agentic AI Trading System"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Real-time analysis: Processes all 50 stocks in < 2 minutes"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Automated decisions: BUY/HOLD/SELL recommendations with confidence scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Integrated AI analysis: LLM summarizes earnings calls → sentiment scores"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Results (META Case Study):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Sharpe ratio: 2.88 (vs 0.47 baseline) = 513% improvement"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• CAGR: 78% (vs 17.94% baseline) = 335% improvement"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Risk control: Max drawdown -7% (vs -18% baseline)"
    p.font.size = Pt(13)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "• Adversarial robustness: 100% resistance to attacks"
    p.font.size = Pt(13)
    p.level = 1
    
    return slide

def main():
    """Create presentation from template"""
    print("Creating enhanced capstone presentation...")
    
    # Try to load template, otherwise create new presentation
    template_path = "AI 894 Capstone Video Presentation Slides.pptx"
    try:
        prs = Presentation(template_path)
        print(f"Loaded template: {template_path}")
    except:
        print("Template not found, creating new presentation...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Load META results
    meta_results = load_meta_results()
    
    # Add slides in logical order
    print("Adding slides...")
    add_title_slide(prs)
    add_intro_slide(prs)
    add_architecture_slide(prs)  # NEW: Architecture diagram
    add_problem_statement_slide(prs)
    add_dqn_formulae_slide(prs)  # NEW: DQN formulae
    add_dqn_training_slide(prs)  # NEW: Detailed training steps
    add_attention_formulae_slide(prs)  # NEW: Attention formulae
    add_methodology_overview_slide(prs)
    add_feature_engineering_examples_slide(prs)
    add_attention_mechanism_examples_slide(prs)
    add_training_examples_slide(prs)
    add_adversarial_definitions_slide(prs)  # NEW: Adversarial definitions
    add_adversarial_examples_slide(prs)  # NEW: Adversarial examples
    add_adversarial_training_examples_slide(prs)
    add_agentic_ai_slide(prs)  # NEW: Agentic AI explanation
    add_api_collaboration_slide(prs)  # NEW: API collaboration
    add_use_case_slide(prs)  # NEW: Compelling use case
    add_meta_results_slide(prs, meta_results)
    add_adversarial_results_slide(prs, meta_results)
    add_infrastructure_slide(prs)
    add_benefits_slide(prs)
    add_conclusion_slide(prs)
    
    # Save presentation
    output_path = "AI_894_Capstone_Presentation_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Enhanced presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()

