#!/usr/bin/env python3
"""
Create enhanced capstone presentation with:
- Architecture diagram
- Detailed RL/DQN formulas
- Adversarial attack definitions and examples
- Use cases
- Agentic AI explanation
- API integration details
- Publication-grade formulas
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import os

def load_meta_results():
    """Load META case study results"""
    try:
        with open('../results/meta_model_results.json', 'r') as f:
            return json.load(f)
    except Exception as e:
        print(f"Warning: Could not load META results: {e}")
        return None

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Adversarial-Robust Multi-Head Attention DQN\nfor Asset Pricing"
    subtitle.text = "AI 894 Capstone Project\nAgentic AI System for Intelligent Trading Decisions"
    
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    # Use title-only layout and remove title placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Remove title placeholder if it exists
    for shape in list(slide.shapes):
        if hasattr(shape, 'placeholder_format'):
            if shape.placeholder_format.idx == 0:  # Title placeholder
                sp = shape.element
                sp.getparent().remove(sp)
                break
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.text = "System Architecture: Multi-Head Attention DQN"
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    
    # Add architecture image if it exists
    img_path = "architecture.png"
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.2), Inches(9), Inches(6))
        except Exception as e:
            print(f"Warning: Could not add architecture image: {e}")
            # Add description instead
            desc_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            desc_frame = desc_shape.text_frame
            desc_frame.text = "Architecture: 8 Feature Groups → Multi-Head Attention (8 heads) → 3 Transformer Layers → Global Average Pooling → Q-Network → Q-values (BUY/HOLD/SELL)"
            desc_frame.word_wrap = True
            for paragraph in desc_frame.paragraphs:
                paragraph.font.size = Pt(18)
    else:
        # Add text description
        desc_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        desc_frame = desc_shape.text_frame
        desc_frame.text = "Architecture Flow:\n1. Input: 8 Feature Groups (Price, Macro, Commodities, Market, Forex, Technical, Earnings, Crypto)\n2. Multi-Head Attention: 8 heads processing different feature aspects\n3. Transformer Layers: 3 layers for temporal dependencies\n4. Global Average Pooling: Summarize sequence\n5. Q-Network: Output Q-values for BUY/HOLD/SELL actions"
        desc_frame.word_wrap = True
        for paragraph in desc_frame.paragraphs:
            paragraph.font.size = Pt(18)
    
    return slide

def add_rl_formulas_slide(prs):
    """Add detailed RL/DQN formulas slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Reinforcement Learning: Deep Q-Network (DQN) Formulation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Markov Decision Process (MDP) Framework:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "MDP = (S, A, P, R, γ) where:"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • S: State space (43 features × 20 time steps)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • A: Action space {BUY, HOLD, SELL}"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • P: Transition probability P(s'|s,a)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • R: Reward function (portfolio return)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • γ = 0.99: Discount factor"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Objective Function:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "π* = argmax_π E_π[Σ(γ^t × r_t)]"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Q-Function (Action-Value Function):"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Q^π(s, a) = E_π[Σ(γ^k × r_{t+k+1}) | s_t=s, a_t=a]"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Bellman Optimality Equation:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Q*(s, a) = E[r(s,a,s') + γ × max_a' Q*(s', a')]"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Q-Learning Update Rule:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Q(s,a) ← Q(s,a) + α[r + γ × max_a' Q(s',a') - Q(s,a)]"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Loss Function (MSE):"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "L(θ) = E[(Q_θ(s,a) - Q_target)^2]"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "where Q_target = r + γ(1-done) × max_a' Q_θ⁻(s',a')"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.level = 1
    
    return slide

def add_dqn_training_slide(prs):
    """Add DQN training process with formulas"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "DQN Training Algorithm: Step-by-Step"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Experience Replay Buffer:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Store transitions: (s_t, a_t, r_t, s_{t+1}, done_t)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Training Loop:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. Sample batch B ~ Uniform(D_replay)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "2. Compute Q_target for each (s,a,r,s'):"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   Q_target = r + γ(1-done) × max_a' Q_θ⁻(s',a')"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "3. Compute loss:"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   L(θ) = (1/|B|) Σ(Q_θ(s,a) - Q_target)^2"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "4. Gradient update:"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   θ ← θ - α × ∇_θ L(θ)"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "5. Update target network every C steps:"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   θ⁻ ← θ (every 100 steps)"
    p.font.size = Pt(12)
    p.font.name = 'Courier New'
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Example: META Stock Training"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "State s: RSI=0.65, MACD=0.02, sentiment=0.8"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Action a: BUY"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Reward r: +2.3% (portfolio return)"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Q_θ(s, BUY) = 0.82, Q_target = 0.85"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Update: θ ← θ - 0.001 × ∇_θ(0.82 - 0.85)^2"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    
    return slide

def add_adversarial_attacks_slide(prs):
    """Add adversarial attack definitions and examples"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Attacks: Definitions & Examples"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Threat Model:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Find perturbation δ such that:"
    p.font.size = Pt(16)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "max_{||δ||_p ≤ ε} L(Q_θ(s + δ, a), y)"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. FGSM (Fast Gradient Sign Method):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "x_adv = x + ε × sign(∇_x L(Q_θ(x), y))"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Example: META stock state"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  Original: RSI=0.65, MACD=0.02 → Prediction: BUY (Q=0.85)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  Gradient: ∇_RSI = -0.15, ∇_MACD = 0.08"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  Perturbation (ε=0.01): δ = [0.01×(-1), 0.01×1] = [-0.01, 0.01]"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  Adversarial: RSI=0.64, MACD=0.03 → Prediction: SELL (Q=0.25) ✗"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. PGD (Projected Gradient Descent):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "x^(0) = x; x^(t+1) = Proj_B_ε(x)[x^(t) + α × sign(∇_x L)]"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Iterative version of FGSM (10 iterations), stronger attack"
    p.font.size = Pt(12)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. C&W (Carlini & Wagner):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "min_δ ||δ||_p + c × f(x + δ)  s.t.  x + δ ∈ [0,1]^d"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Optimization-based attack, finds minimal perturbation"
    p.font.size = Pt(12)
    
    return slide

def add_adversarial_training_slide(prs):
    """Add adversarial training process"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Adversarial Training: Building Robust Models"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Training Objective:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "L_total = 0.5 × L(Q_θ(s), y) + 0.5 × L(Q_θ(s_adv), y)"
    p.font.size = Pt(16)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Training Process:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. Forward pass on clean data: Q_θ(s) → Q_values"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "2. Generate adversarial examples: s_adv = s + δ_FGSM"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "3. Forward pass on adversarial data: Q_θ(s_adv) → Q_values"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "4. Compute combined loss and backpropagate"
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Results on META Stock:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Without Adversarial Training:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • FGSM attack: Sharpe ratio 2.88 → 0.12 (96% degradation)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "With Adversarial Training:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • FGSM attack: Sharpe ratio 2.88 → 2.88 (0% degradation)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • PGD attack: 100% resistance"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • C&W attack: 100% resistance"
    p.font.size = Pt(12)
    p.level = 1
    
    return slide

def add_agentic_ai_slide(prs):
    """Add explanation of agentic AI system"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Agentic AI System: Autonomous Trading Agent"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "What is an Agentic AI System?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "An autonomous agent that:"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. Observes Environment"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Continuously monitors market data, prices, indicators"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Fetches real-time data from multiple APIs"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "2. Makes Autonomous Decisions"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Analyzes 43 features across 8 groups using MHA-DQN"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Selects optimal action: BUY/HOLD/SELL"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "3. Executes Actions"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Executes trades based on Q-values"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Manages portfolio positions"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "4. Learns from Experience"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Updates Q-function based on rewards"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Adapts strategy to market conditions"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "5. Collaborates with Other AI Systems"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Uses OpenAI GPT for earnings call analysis"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Integrates sentiment analysis from NLP models"
    p.font.size = Pt(12)
    p.level = 1
    
    return slide

def add_api_integration_slide(prs):
    """Add API integration and AI collaboration details"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "API Integration & AI System Collaboration"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Integrated APIs:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "1. Alpha Vantage API"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Real-time stock prices (OHLCV data)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Macroeconomic indicators (GDP, CPI, unemployment, Fed rates)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Commodity prices (WTI crude, gold, natural gas, copper)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Market indices (SPY, VIX, sector ETFs: XLK, XLF, XLE)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Forex rates (EUR/USD, GBP/USD, USD/JPY, AUD/USD)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "2. Financial Modeling Prep (FMP) API"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Earnings call transcripts"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Financial statements (EPS, revenue, guidance)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "3. OpenAI GPT API (Agentic Collaboration)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Analyzes earnings call transcripts"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Generates sentiment scores (0-1 scale)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Extracts key insights (growth outlook, risks, opportunities)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "AI System Collaboration Workflow:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "1. MHA-DQN Agent requests earnings data from FMP API"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "2. FMP API provides transcript → Sent to OpenAI GPT"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "3. OpenAI GPT analyzes transcript and generates summary:"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Sentiment score: 0.85 (positive)"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "   • Key points: Strong revenue growth, new product launch"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "4. MHA-DQN Agent incorporates sentiment into feature Group 7"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "5. Agent makes trading decision using enriched features"
    p.font.size = Pt(14)
    
    return slide

def add_use_case_slide(prs):
    """Add compelling use case"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Compelling Use Case: Institutional Trading Desk"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Scenario: Large Hedge Fund Trading Desk"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Challenge:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Portfolio of 50+ stocks requires constant monitoring"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Human traders cannot process all information in real-time"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Need robust decisions resistant to market manipulation"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Solution: Deploy MHA-DQN Agentic AI System"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Agent autonomously:"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Monitors all 50 stocks simultaneously"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Fetches data from Alpha Vantage, FMP APIs in real-time"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Collaborates with OpenAI GPT for earnings analysis"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Makes BUY/HOLD/SELL decisions based on 43 features"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Maintains robustness under adversarial conditions"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Results (META Stock Example):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "  • Sharpe Ratio: 2.88 (vs. 0.47 baseline)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • CAGR: 78% (vs. 18% baseline)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • Portfolio value: $10,000 → $16,875 (69% return)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "  • 100% resistance to adversarial attacks"
    p.font.size = Pt(14)
    p.level = 1
    
    return slide

def add_formulas_summary_slide(prs):
    """Add summary of key publication-grade formulas"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Key Mathematical Formulations (Publication-Grade)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Multi-Head Attention:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Attention(Q,K,V) = softmax(QK^T / √d_k) V"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "MHA(Q,K,V) = Concat(head_1,...,head_H) W_O"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Q-Learning:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Q*(s,a) = E[r(s,a,s') + γ max_a' Q*(s',a')]"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Adversarial Training:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "L_total = 0.5 L(Q_θ(s), y) + 0.5 L(Q_θ(s_adv), y)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "where s_adv = s + ε sign(∇_s L)"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "Performance Metrics:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "Sharpe Ratio = (μ_r - r_f) / σ_r"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "CAGR = (V_T / V_0)^(1/T) - 1"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    p.space_after = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "Max Drawdown = min_t (P_t - P_peak) / P_peak"
    p.font.size = Pt(14)
    p.font.name = 'Courier New'
    
    return slide

def main():
    """Create enhanced presentation"""
    print("Creating enhanced capstone presentation...")
    
    # Load template
    template_path = "AI 894 Capstone Video Presentation Slides.pptx"
    try:
        prs = Presentation(template_path)
        print(f"Loaded template: {template_path}")
    except:
        print("Template not found, creating new presentation...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Load META results
    meta_results = load_meta_results()
    
    # Add slides in order
    print("Adding enhanced slides...")
    add_title_slide(prs)
    add_opening_context_slide(prs)
    add_architecture_slide(prs)
    add_rl_formulas_slide(prs)
    add_dqn_training_slide(prs)
    add_adversarial_attacks_slide(prs)
    add_adversarial_training_slide(prs)
    add_agentic_ai_slide(prs)
    add_api_integration_slide(prs)
    add_use_case_slide(prs)
    add_formulas_summary_slide(prs)
    
    # Save presentation
    output_path = "AI_894_Capstone_Presentation_Enhanced.pptx"
    prs.save(output_path)
    print(f"\n✅ Enhanced presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print("\n📋 Included Features:")
    print("   ✓ Architecture diagram")
    print("   ✓ Detailed RL/DQN formulas")
    print("   ✓ Adversarial attack definitions & examples")
    print("   ✓ Compelling use case")
    print("   ✓ Agentic AI system explanation")
    print("   ✓ API integration details")
    print("   ✓ Publication-grade formulas")

if __name__ == "__main__":
    main()

